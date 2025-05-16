# Copyright 2024 Liu Siyao
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

from __future__ import print_function

import logging
import os
from functools import partial
from operator import attrgetter
from typing import List, Union
from pathlib import Path
import io # Add io import

from PIL import Image
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from rapidfuzz import process as fuze_process
from tqdm import tqdm
from pptx.util import Emu # For EMU to Px conversion
from pptx.shapes.picture import Picture # Added for type hinting

from pptx2md.multi_column import get_multi_column_slide_if_present
from pptx2md.types import (
    ConversionConfig,
    GeneralSlide,
    ImageElement,
    ListItemElement,
    ParagraphElement,
    ParsedPresentation,
    SlideElement,
    TableElement,
    TextRun,
    TextStyle,
    TitleElement,
    CodeBlockElement,
    FormulaElement,
)
from pptx2md.utils import emu_to_px # Assuming emu_to_px is now in utils

logger = logging.getLogger(__name__)

picture_count = 0

def is_code_font(font) -> bool:
    """Checks if the font is a common monospaced/code font."""
    if font and font.name:
        # Common monospaced fonts; .lower() for case-insensitivity.
        # Add more as needed.
        monospaced_fonts = ["consolas", "courier new", "menlo", "menlo regular", "monaco", "lucida console", "dejavu sans mono"]
        if font.name.lower() in monospaced_fonts:
            return True
    return False


def is_title(shape):
    if shape.is_placeholder and (shape.placeholder_format.type == PP_PLACEHOLDER.TITLE or
                                 shape.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE or
                                 shape.placeholder_format.type == PP_PLACEHOLDER.VERTICAL_TITLE or
                                 shape.placeholder_format.type == PP_PLACEHOLDER.CENTER_TITLE):
        return True
    return False


def is_text_block(config: ConversionConfig, shape):
    if shape.has_text_frame:
        if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.BODY:
            return True
        if len(shape.text) > config.min_block_size:
            return True
    return False


def is_list_block(shape) -> bool:
    levels = []
    for para in shape.text_frame.paragraphs:
        if para.level not in levels:
            levels.append(para.level)
        if para.level != 0 or len(levels) > 1:
            return True
    return False


def is_accent(font):
    return font and (
        font.underline or font.italic or (
            font.color.type == MSO_COLOR_TYPE.SCHEME and
        (font.color.theme_color == MSO_THEME_COLOR.ACCENT_1 or 
         font.color.theme_color == MSO_THEME_COLOR.ACCENT_2 or
         font.color.theme_color == MSO_THEME_COLOR.ACCENT_3 or 
         font.color.theme_color == MSO_THEME_COLOR.ACCENT_4 or
         font.color.theme_color == MSO_THEME_COLOR.ACCENT_5 or 
         font.color.theme_color == MSO_THEME_COLOR.ACCENT_6))
        )

def is_math(text):
    return text and ( text.startswith('$') and text.endswith('$') )

def is_strong(font):
    return font and (font.bold or 
                     (font.color.type == MSO_COLOR_TYPE.SCHEME and 
                      (font.color.theme_color == MSO_THEME_COLOR.DARK_1 or 
                       font.color.theme_color == MSO_THEME_COLOR.DARK_2)))


def get_text_runs(para) -> List[TextRun]:
    runs = []
    for run in para.runs:
        # Populate font_name directly in TextRun
        font_name = run.font.name if run.font and hasattr(run.font, 'name') else None
        result = TextRun(text=run.text, style=TextStyle(), font_name=font_name)
        if result.text == '':
            continue
        try:
            result.style.hyperlink = run.hyperlink.address
        except:
            pass
            #result.style.hyperlink = 'error:ppt-link-parsing-issue'
        if is_accent(run.font):
            result.style.is_accent = True
        if is_strong(run.font):
            result.style.is_strong = True
        if is_math(run.text):
            result.style.is_math = True
        if run.font and (run.font.color.type == MSO_COLOR_TYPE.RGB):
            result.style.color_rgb = run.font.color.rgb
        if run.font and is_code_font(run.font):
            result.style.is_code = True
        runs.append(result)
    return runs


def process_title(config: ConversionConfig, shape, slide_idx) -> TitleElement:
    text = shape.text_frame.text.strip()
    if config.custom_titles:
        res = fuze_process.extractOne(text, config.custom_titles.keys(), score_cutoff=92)
        if not res:
            return TitleElement(content=text.strip(), level=max(config.custom_titles.values()) + 1)
        else:
            logger.info(f'Title in slide {slide_idx} "{text}" is converted to "{res[0]}" as specified in title file.')
            return TitleElement(content=res[0].strip(), level=config.custom_titles[res[0]])
    else:
        return TitleElement(content=text.strip(), level=1)


def process_text_blocks(config: ConversionConfig, shape, slide_idx) -> List[Union[ListItemElement, ParagraphElement]]:
    results = []
    
    # if slide_idx == 52: breakpoint()

    if not shape.has_text_frame or not shape.text_frame.paragraphs:
        return results

    # Determine if the shape's substantive content is entirely code-styled.
    all_substantive_content_is_code = True
    has_any_substantive_content = False

    for para in shape.text_frame.paragraphs:
        # We need to get runs with their styles to check .is_code
        para_runs = get_text_runs(para)
        if not para_runs and not para.text.strip(): # Paragraph is empty or only whitespace
            # This paragraph contributes to structure but not to substantive content check
            # unless it's the *only* thing in the shape.
            # We'll add an empty code run for it later if shape is deemed code.
            continue 

        for run in para_runs:
            if run.text.strip(): # Substantive text
                has_any_substantive_content = True
                if not run.style.is_code:
                    all_substantive_content_is_code = False
                    break # Found non-code substantive text
            # Non-substantive runs (whitespace) don't break the "all code" criteria
        if not all_substantive_content_is_code:
            break
    
    # If there's no substantive content at all, it cannot be "all code" in a meaningful way.
    if not has_any_substantive_content:
        all_substantive_content_is_code = False

    if all_substantive_content_is_code:
        # This shape's substantive content is entirely code-styled.
        # Process all its paragraphs as ParagraphElements (for potential merging into CodeBlockElement).
        for para in shape.text_frame.paragraphs:
            para_runs = get_text_runs(para) # Re-get runs with style
            # If paragraph is empty (no runs after get_text_runs, e.g. only filtered chars, or truly empty para.text),
            # create a ParagraphElement with a single empty, code-styled TextRun to preserve the line.
            if not para_runs and not para.text: # Truly empty line
                 results.append(ParagraphElement(content=[TextRun(text="", style=TextStyle(is_code=True))]))
            elif para_runs: # Paragraph has runs (could be actual code, or just whitespace forming a run)
                 # Ensure all runs within this paragraph are also marked as code if the whole shape is code.
                 # This might be redundant if get_text_runs already did it perfectly, but ensures consistency here.
                updated_para_runs = [
                    TextRun(text=r.text, style=TextStyle(is_code=True, 
                                                        hyperlink=r.style.hyperlink, # Preserve links if any
                                                        color_rgb=r.style.color_rgb)) # Preserve color if any
                    for r in para_runs
                ]
                results.append(ParagraphElement(content=updated_para_runs))
            # else: para.text might have content but get_text_runs yielded nothing (e.g. filtered chars).
            # This case should ideally not happen if get_text_runs is robust.
            # If it does, such a paragraph is currently skipped.
            # To ensure an empty line, we can add:
            elif not para_runs and para.text.isspace(): # Whitespace only line
                results.append(ParagraphElement(content=[TextRun(text=para.text, style=TextStyle(is_code=True))]))

    else:
        # Not an entirely code-styled shape. Fall back to list/paragraph detection.
        is_list = is_list_block(shape)
        for para in shape.text_frame.paragraphs:
            # For non-code content, skip paragraphs that are visually empty.
            if not para.text.strip(): 
                continue
            
            text_runs = get_text_runs(para)
            if not text_runs: # If, after processing, text_runs is empty
                continue

            if is_list:
                results.append(ListItemElement(content=text_runs, level=para.level))
            else:
                results.append(ParagraphElement(content=text_runs))
                
    return results


def _crop_image_if_needed(
    img_to_process: Image.Image, 
    crop_l_pct: float, 
    crop_r_pct: float, 
    crop_t_pct: float, 
    crop_b_pct: float,
    current_pil_format: str,
    original_blob: bytes,
    slide_idx: int,
    config: ConversionConfig
) -> tuple[Union[Image.Image, None], bytes, Union[tuple[int, int], None], tuple[Union[float, None], Union[float, None], Union[float, None], Union[float, None]]]:
    """
    Applies cropping to a Pillow Image object if specified and enabled.
    Returns the (potentially) cropped image object, its blob, its new dimensions,
    and the crop percentages to be stored in ImageElement.
    """
    pil_original_w, pil_original_h = img_to_process.size
    cropped_img_obj = img_to_process
    current_image_blob = original_blob # Start with original if crop fails or not applied
    final_blob_w_px, final_blob_h_px = pil_original_w, pil_original_h
    
    crop_l_for_element, crop_r_for_element, crop_t_for_element, crop_b_for_element = None, None, None, None

    has_crop_info = (crop_l_pct > 0.00001 or crop_r_pct > 0.00001 or 
                     crop_t_pct > 0.00001 or crop_b_pct > 0.00001)

    if has_crop_info and config.apply_cropping_in_parser:
        left = int(round(pil_original_w * crop_l_pct))
        top = int(round(pil_original_h * crop_t_pct))
        right = int(round(pil_original_w * (1.0 - crop_r_pct)))
        bottom = int(round(pil_original_h * (1.0 - crop_b_pct)))

        if left < right and top < bottom:
            try:
                cropped_img_obj = img_to_process.crop((left, top, right, bottom))
                logger.info(f'Image in slide {slide_idx} pre-cropped in parser. New blob dims: {cropped_img_obj.size}')
                with io.BytesIO() as cropped_blob_io:
                    save_format = current_pil_format if current_pil_format else 'PNG'
                    try:
                        cropped_img_obj.save(cropped_blob_io, format=save_format)
                    except KeyError: 
                        logger.warning(f"Format {save_format} not supported by Pillow for saving, falling back to PNG.")
                        save_format = 'PNG'
                        if cropped_img_obj.mode != 'RGBA' and cropped_img_obj.mode != 'RGB':
                            cropped_img_obj = cropped_img_obj.convert('RGBA')
                        cropped_img_obj.save(cropped_blob_io, format=save_format)
                    current_image_blob = cropped_blob_io.getvalue()
                
                final_blob_w_px, final_blob_h_px = cropped_img_obj.size
                # Cropping applied, so no percentages needed for ImageElement
                crop_l_for_element, crop_r_for_element, crop_t_for_element, crop_b_for_element = None, None, None, None
            except Exception as e:
                logger.warning(f"Failed to apply crop in parser for image in slide {slide_idx}. Error: {e}")
                # Fallback: use uncropped dimensions, and pass crop info
                cropped_img_obj = img_to_process # Revert to original if crop failed
                current_image_blob = original_blob # Revert to original blob
                final_blob_w_px, final_blob_h_px = pil_original_w, pil_original_h
                if has_crop_info: # Still pass the crop info if it existed
                    crop_l_for_element, crop_r_for_element, crop_t_for_element, crop_b_for_element = \
                        crop_l_pct, crop_r_pct, crop_t_pct, crop_b_pct
        else:
            logger.warning(f"Invalid crop dimensions for image in slide {slide_idx}, not applying crop in parser.")
            final_blob_w_px, final_blob_h_px = pil_original_w, pil_original_h
            if has_crop_info: # Pass the crop info
                crop_l_for_element, crop_r_for_element, crop_t_for_element, crop_b_for_element = \
                    crop_l_pct, crop_r_pct, crop_t_pct, crop_b_pct
    
    elif has_crop_info: # Cropping info exists, but apply_cropping_in_parser is False
        logger.info(f"Parser cropping disabled for image in slide {slide_idx}. Crop info will be passed to formatter.")
        final_blob_w_px, final_blob_h_px = pil_original_w, pil_original_h # Dimensions of uncropped blob
        crop_l_for_element, crop_r_for_element, crop_t_for_element, crop_b_for_element = \
            crop_l_pct, crop_r_pct, crop_t_pct, crop_b_pct
    else: # No crop defined on shape, or Pillow couldn't open
        # final_blob_w_px, final_blob_h_px already set from img_to_process.size
        crop_l_for_element, crop_r_for_element, crop_t_for_element, crop_b_for_element = None, None, None, None

    return cropped_img_obj, current_image_blob, (final_blob_w_px, final_blob_h_px), \
           (crop_l_for_element, crop_r_for_element, crop_t_for_element, crop_b_for_element)


def _handle_wmf_conversion(
    current_image_blob: bytes,
    original_pic_ext: str,
    config: ConversionConfig,
    slide_idx: int,
    picture_idx_for_temp_name: int
) -> tuple[bytes, str, str, bool]:
    """
    Handles WMF to PNG conversion if WMF is detected and not disabled.
    Returns the image blob, pic_ext, pil_format, and a flag indicating if conversion occurred.
    """
    updated_blob = current_image_blob
    updated_ext = original_pic_ext
    updated_pil_format = pil_format_map.get(original_pic_ext, original_pic_ext.upper()) # Calculate initial PIL format
    converted_from_wmf = False

    if original_pic_ext == 'wmf':
        if config.disable_wmf:
            logger.info(f"WMF image in slide {slide_idx} will be kept as original WMF (processing disabled).")
            # Blob, ext, and pil_format remain as is for WMF
        else:
            try:
                temp_wmf_path = Path(config.image_dir) / f"__temp_wmf_{picture_idx_for_temp_name}.wmf"
                with open(temp_wmf_path, 'wb') as tmp_f:
                    tmp_f.write(current_image_blob)

                from wand.image import Image as WandImage
                with WandImage(filename=str(temp_wmf_path)) as img:
                    img.format = 'png'
                    with io.BytesIO() as png_blob_io:
                        img.save(file=png_blob_io)
                        updated_blob = png_blob_io.getvalue()
                
                updated_pil_format = 'PNG'
                updated_ext = 'png'
                converted_from_wmf = True
                logger.info(f'WMF image in slide {slide_idx} converted to PNG for processing.')
                
                if temp_wmf_path.exists():
                    try:
                        os.remove(temp_wmf_path)
                    except OSError as e:
                        logger.warning(f"Could not remove temp WMF file {temp_wmf_path}: {e}")
            except Exception as e:
                logger.warning(
                    f'Cannot convert WMF image in slide {slide_idx} to PNG. Error: {e}. '
                    f'Attempting to use original WMF if possible.'
                )
                # Fallback: blob, ext, pil_format remain as original WMF
                updated_blob = current_image_blob 
                updated_ext = original_pic_ext
                updated_pil_format = pil_format_map.get(original_pic_ext, original_pic_ext.upper())

    return updated_blob, updated_ext, updated_pil_format, converted_from_wmf


def _handle_tiff_conversion(
    current_image_blob: bytes,
    original_pic_ext: str, # This ext might be 'png' if WMF conversion happened
    current_pil_format: str,
    slide_idx: int
) -> tuple[bytes, str, str]:
    """
    Handles TIFF to PNG conversion if TIFF is detected.
    Assumes WMF conversion (if any) has already occurred.
    Returns the image blob, pic_ext, and pil_format.
    """
    updated_blob = current_image_blob
    updated_ext = original_pic_ext
    updated_pil_format_out = current_pil_format

    # Check original_pic_ext which is the true source before this stage (or after WMF conversion)
    if original_pic_ext in ['tif', 'tiff']:
        logger.info(f"Image in slide {slide_idx} is TIFF ({original_pic_ext}), attempting conversion to PNG.")
        try:
            img_from_tiff: Image.Image
            with io.BytesIO(current_image_blob) as tiff_bytes_io:
                img_from_tiff = Image.open(tiff_bytes_io)
                img_from_tiff.load() 
            
            if img_from_tiff.mode != 'RGBA':
                img_from_tiff = img_from_tiff.convert('RGBA')

            with io.BytesIO() as png_bytes_io:
                img_from_tiff.save(png_bytes_io, format='PNG')
                updated_blob = png_bytes_io.getvalue()
            
            updated_ext = 'png'
            updated_pil_format_out = 'PNG'
            logger.info(f"TIFF image in slide {slide_idx} successfully converted to PNG.")
        except Exception as e:
            logger.warning(
                f"Failed to convert TIFF image in slide {slide_idx} from .{original_pic_ext} to PNG. Error: {e}. "
                f"Proceeding with original TIFF data. Markdown output may not display this image."
            )
            # Fallback: blob, ext, pil_format remain as original TIFF
            updated_blob = current_image_blob
            updated_ext = original_pic_ext
            updated_pil_format_out = pil_format_map.get(original_pic_ext, original_pic_ext.upper())


    return updated_blob, updated_ext, updated_pil_format_out


def _open_and_prepare_image_with_pillow(
    image_blob: bytes,
    pic_ext: str, # Extension of the blob (could be original, or png after wmf/tiff conversion)
    pil_format_for_opening: str, # PIL format string for the blob
    slide_idx: int,
    shape, # Original pptx shape, for fallback dimensions
    is_wmf_processing_disabled_for_this_image: bool 
) -> tuple[Union[Image.Image, None], Union[int, None], Union[int, None]]:
    """
    Opens an image blob with Pillow, prepares it (e.g., converts mode), 
    and returns the Pillow image object and its dimensions.
    Handles cases where Pillow cannot open the image or WMF processing is disabled.
    """
    img_to_process = None
    blob_w_px, blob_h_px = None, None

    # Skip Pillow processing entirely for WMFs if config.disable_wmf is True.
    # The pic_ext here would still be 'wmf'.
    if pic_ext == 'wmf' and is_wmf_processing_disabled_for_this_image:
        logger.info(f"Pillow processing skipped for WMF image in slide {slide_idx} as per config.")
        if hasattr(shape.image, 'size') and shape.image.size:
            blob_w_px, blob_h_px = shape.image.size
        return None, blob_w_px, blob_h_px

    try:
        img_to_process = Image.open(io.BytesIO(image_blob))
        # Ensure image data is loaded to prevent issues with closed streams later if blob was from BytesIO
        img_to_process.load()

        # Standardize modes for common output formats or problematic input modes
        if pil_format_for_opening == 'PNG' and img_to_process.mode != 'RGBA':
            img_to_process = img_to_process.convert('RGBA')
        elif img_to_process.mode == 'P': # Paletted images
            # Convert to RGBA if transparency might be present, else RGB
            img_to_process = img_to_process.convert('RGBA' if 'A' in img_to_process.info.get('transparency', []) else 'RGB')
        elif img_to_process.mode == 'CMYK':
            img_to_process = img_to_process.convert('RGB')
        
        blob_w_px, blob_h_px = img_to_process.size
    except Exception as e:
        logger.warning(
            f"Pillow could not open/process image blob for slide {slide_idx} (ext: {pic_ext}, PIL format: {pil_format_for_opening}). Error: {e}. "
            f"Skipping Pillow processing for this image."
        )
        img_to_process = None
        if hasattr(shape.image, 'size') and shape.image.size: # Fallback to shape.image.size
            blob_w_px, blob_h_px = shape.image.size
        # else blob_w_px, blob_h_px remain None

    return img_to_process, blob_w_px, blob_h_px


def _save_image_and_get_path(
    config: 'ConversionConfig',
    image_blob: bytes,
    image_final_ext: str,
    original_shape_id_for_naming: Union[str, int],
    slide_id: int
) -> str:
    """
    Saves an image (using the provided blob and extension) and returns its relative path.
    The image is saved to a subdirectory, structured as:
    <actual_images_base_dir>/<config.pptx_path.stem>_img/<generated_filename>
    The returned path is relative to actual_images_base_dir.
    It attempts to derive actual_images_base_dir correctly if config.output_dir seems to be a file path.
    """
    if not config.output_dir:
        raise ValueError("config.output_dir must be set to save images.")

    actual_images_base_dir = config.output_dir
    # Attempt to determine the correct base directory for images
    # if config.output_dir appears to be a full file path.
    if not config.output_dir.is_dir():
        if config.output_dir.suffix:  # Has a file-like extension (e.g., .md)
            logger.info(
                f"config.output_dir ('{config.output_dir}') appears to be a full file path. "
                f"Using its parent directory ('{config.output_dir.parent}') as the base for the images folder."
            )
            actual_images_base_dir = config.output_dir.parent
        else:
            # Not a directory and no suffix, could be problematic.
            logger.warning(
                f"config.output_dir ('{config.output_dir}') is not recognized as a directory and lacks a file extension. "
                f"Attempting to use it as the base for the images folder. This may fail."
            )
    
    # Ensure the determined base directory for images (e.g., "outputs/") exists.
    actual_images_base_dir.mkdir(parents=True, exist_ok=True)

    if not config.pptx_path or not config.pptx_path.stem:
        raise ValueError("config.pptx_path.stem must be available for naming the image subfolder.")

    # New image subfolder name: <config.pptx_path.stem>_img
    image_subfolder_name = f"{config.pptx_path.stem}_img"
    
    # Define the specific target directory for this presentation's images
    # e.g., <actual_images_base_dir>/presentation_name_img/
    image_target_dir = actual_images_base_dir / image_subfolder_name
    
    # Create this specific image directory
    image_target_dir.mkdir(parents=True, exist_ok=True)

    filename_prefix = f"slide{slide_id}_shape{original_shape_id_for_naming}"
    filename = f"{filename_prefix}.{image_final_ext}" # Use the final, correct extension

    saved_image_full_path = image_target_dir / filename

    with open(saved_image_full_path, "wb") as f:
        f.write(image_blob) # Save the processed image blob
    
    # Path should be relative to 'actual_images_base_dir'
    # e.g., if actual_images_base_dir is 'outputs', path is 'presentation_name_img/file.png'
    relative_path = Path(image_subfolder_name) / filename
    
    return str(relative_path)


# pil_format_map needs to be accessible by the helper functions if they are top-level
# or passed as an argument. For now, keeping it module-level.
pil_format_map = {'jpg': 'JPEG', 'jpeg': 'JPEG', 'tif': 'TIFF', 'tiff': 'TIFF'}


def process_picture(config: ConversionConfig, shape: Picture, slide_idx: int) -> Union[ImageElement, None]:
    if config.disable_image:
        return None

    if not hasattr(shape, 'image') or not shape.image:
        logger.warning(f"Shape in slide {slide_idx} seems to be a picture but has no image data, skipped.")
        return None

    global picture_count # Used for unique naming and WMF temp file

    # Initial properties from shape.image
    initial_pic_ext = shape.image.ext.lower()
    current_image_blob = shape.image.blob
    # current_pil_format will be determined after potential WMF/TIFF conversions.
    
    # WMF Conversion
    current_image_blob, effective_pic_ext, effective_pil_format, converted_from_wmf = \
        _handle_wmf_conversion(current_image_blob, initial_pic_ext, config, slide_idx, picture_count)

    # TIFF to PNG Conversion (if not already PNG from WMF)
    if not converted_from_wmf and effective_pic_ext in ['tif', 'tiff']:
        current_image_blob, effective_pic_ext, effective_pil_format = \
            _handle_tiff_conversion(current_image_blob, effective_pic_ext, effective_pil_format, slide_idx)
    
    # Image Cropping and Pillow Processing
    img_to_process_for_crop = None
    final_blob_w_px, final_blob_h_px = None, None 
    crop_l_for_element, crop_r_for_element, crop_t_for_element, crop_b_for_element = None, None, None, None

    is_wmf_and_disabled = (initial_pic_ext == 'wmf' and config.disable_wmf)

    img_to_process_for_crop, initial_blob_w_px, initial_blob_h_px = \
        _open_and_prepare_image_with_pillow(
            current_image_blob, 
            effective_pic_ext, 
            effective_pil_format, 
            slide_idx, 
            shape,
            is_wmf_and_disabled
        )
    
    final_blob_w_px, final_blob_h_px = initial_blob_w_px, initial_blob_h_px 

    if img_to_process_for_crop: 
        crop_l_pct = getattr(shape, 'crop_left', 0.0) 
        crop_r_pct = getattr(shape, 'crop_right', 0.0)
        crop_t_pct = getattr(shape, 'crop_top', 0.0)
        crop_b_pct = getattr(shape, 'crop_bottom', 0.0)

        _cropped_img_obj, current_image_blob, \
        (final_blob_w_px, final_blob_h_px), \
        (crop_l_for_element, crop_r_for_element, crop_t_for_element, crop_b_for_element) = \
            _crop_image_if_needed(
                img_to_process_for_crop, 
                crop_l_pct, crop_r_pct, crop_t_pct, crop_b_pct,
                effective_pil_format, 
                current_image_blob,   
                slide_idx,
                config
            )
    
    elif initial_blob_w_px is not None: 
        has_crop_info_on_shape = (getattr(shape, 'crop_left', 0.0) > 0.00001 or
                                  getattr(shape, 'crop_right', 0.0) > 0.00001 or
                                  getattr(shape, 'crop_top', 0.0) > 0.00001 or
                                  getattr(shape, 'crop_bottom', 0.0) > 0.00001)
        if has_crop_info_on_shape:
            crop_l_for_element = getattr(shape, 'crop_left', 0.0)
            crop_r_for_element = getattr(shape, 'crop_right', 0.0)
            crop_t_for_element = getattr(shape, 'crop_top', 0.0)
            crop_b_for_element = getattr(shape, 'crop_bottom', 0.0)
            logger.info(f"Image in slide {slide_idx} (ext: {effective_pic_ext}) could not be opened by Pillow. "
                        f"Crop info from shape will be passed to formatter if present.")

    # Saving the final (potentially processed) blob
    saved_path_str = _save_image_and_get_path(
        config,
        current_image_blob,      # Pass the processed blob
        effective_pic_ext,       # Pass the final extension (e.g., 'png' after conversion)
        shape.shape_id,          # Pass shape_id for consistent naming
        slide_idx                # Pass slide_id
    )
    picture_count += 1 

    # Create ImageElement
    image_data = ImageElement(
        path=saved_path_str,
        original_width_px=final_blob_w_px,   # Dimensions of the saved file blob
        original_height_px=final_blob_h_px,  # Dimensions of the saved file blob
        original_filename=shape.image.filename if hasattr(shape.image, 'filename') else None,
        display_width_px=emu_to_px(shape.width),
        display_height_px=emu_to_px(shape.height),
        left_px=emu_to_px(shape.left),
        top_px=emu_to_px(shape.top),
        rotation=shape.rotation if hasattr(shape, 'rotation') else 0.0,
        crop_left_pct=crop_l_for_element, 
        crop_right_pct=crop_r_for_element,
        crop_top_pct=crop_t_for_element,
        crop_bottom_pct=crop_b_for_element,
        alt_text=shape.name if hasattr(shape, 'name') else None
    )
    return image_data


def process_table(config: ConversionConfig, shape, slide_idx) -> Union[TableElement, None]:
    table = [[sum([get_text_runs(p)
                   for p in cell.text_frame.paragraphs], [])
              for cell in row.cells]
             for row in shape.table.rows]
    if len(table) > 0:
        return TableElement(content=table)
    return None


def ungroup_shapes(shapes) -> List[SlideElement]:
    res = []
    for shape in shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                res.extend(ungroup_shapes(shape.shapes))
            else:
                res.append(shape)
        except Exception as e:
            logger.warning(f'failed to load shape {shape}, skipped. error: {e}')
    return res


def _refine_elements(initial_elements: List[SlideElement], slide_idx: int) -> List[SlideElement]:
    refined_elements: List[SlideElement] = []
    i = 0
    while i < len(initial_elements):
        current_element = initial_elements[i]

        # 1. Check for Math Paragraphs to convert to FormulaElement
        if isinstance(current_element, ParagraphElement) and \
           len(current_element.content) == 1 and \
           current_element.content[0].style.is_math:
            
            math_run = current_element.content[0]
            raw_text = math_run.text.strip() # Strip to handle potential whitespace around $
            formula_content_for_element: str

            if raw_text.startswith('$$') and raw_text.endswith('$$') and len(raw_text) >= 4:
                formula_content_for_element = raw_text[2:-2]
            elif raw_text.startswith('$') and raw_text.endswith('$') and len(raw_text) >= 2:
                formula_content_for_element = raw_text[1:-1]
            else:
                # This case implies is_math was true, but format is unexpected.
                # Log warning and treat as simple text to avoid formatter errors.
                logger.warning(
                    f"Slide {slide_idx}: Math run '{math_run.text}' marked as math "
                    f"but not in expected $...$ or $$...$$ format. Storing raw."
                )
                formula_content_for_element = math_run.text # Store original text
            
            formula_el = FormulaElement(content=formula_content_for_element, position=current_element.position)
            refined_elements.append(formula_el)
            i += 1
            continue

        # 2. Check for Code Blocks (merging consecutive code paragraphs)
        elif isinstance(current_element, ParagraphElement) and \
             current_element.content and \
             all(run.style.is_code for run in current_element.content):
            
            consecutive_code_paras: List[ParagraphElement] = []
            scan_idx = i
            while scan_idx < len(initial_elements):
                element_to_check = initial_elements[scan_idx]
                if isinstance(element_to_check, ParagraphElement) and \
                   element_to_check.content and \
                   all(run.style.is_code for run in element_to_check.content):
                    # Additional check: ensure this paragraph wasn't meant to be a math block
                    # that happened to use a code font and somehow wasn't converted above.
                    # This is a heuristic: if it also looks like a math block, prioritize math.
                    # (This check might be redundant if math conversion is robust)
                    if not (len(element_to_check.content) == 1 and element_to_check.content[0].style.is_math):
                        consecutive_code_paras.append(element_to_check)
                    else: # It's a code-styled math paragraph; stop collecting for this code block
                        break 
                else: # Not a code paragraph, or different element type
                    break
                scan_idx += 1
            
            if not consecutive_code_paras: # Should not happen if outer 'if' was true, but defensive
                refined_elements.append(current_element) # Add current element as is
                i += 1
                continue

            if len(consecutive_code_paras) == 1:
                single_para_element = consecutive_code_paras[0]
                raw_text_content = "".join(run.text for run in single_para_element.content)
                
                if '\n' not in raw_text_content.strip(): # Single line of code
                    refined_elements.append(single_para_element) 
                else: # Single paragraph, but multi-line text: treat as a CodeBlock
                    code_block = CodeBlockElement(content=raw_text_content, 
                                                  position=single_para_element.position,
                                                  language=None) 
                    refined_elements.append(code_block)
                i += 1 # Processed one element
            
            elif len(consecutive_code_paras) > 1: # Multiple consecutive code paragraphs
                code_lines_texts = ["".join(run.text for run in para.content) for para in consecutive_code_paras]
                full_code_content = "\n".join(code_lines_texts)
                first_para_position = consecutive_code_paras[0].position
                
                code_block = CodeBlockElement(content=full_code_content, 
                                              position=first_para_position,
                                              language=None) 
                refined_elements.append(code_block)
                i += len(consecutive_code_paras) # Advance by the number of merged paragraphs
            
        else: # Not a special math paragraph, nor the start of a code paragraph sequence
            refined_elements.append(current_element)
            i += 1
            
    return refined_elements


def process_shapes(config: ConversionConfig, current_shapes, slide_id: int) -> List[SlideElement]:
    initial_elements: List[SlideElement] = []
    for shape in current_shapes:
        if is_title(shape):
            initial_elements.append(process_title(config, shape, slide_id))
        elif is_text_block(config, shape):
            initial_elements.extend(process_text_blocks(config, shape, slide_id))
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                pic = process_picture(config, shape, slide_id)
                if pic:
                    initial_elements.append(pic)
            except AttributeError as e:
                logger.warning(f'Failed to process picture in slide {slide_id}, skipped: {e}')
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = process_table(config, shape, slide_id)
            if table:
                initial_elements.append(table)
        else:
            try:
                ph = shape.placeholder_format
                if ph.type == PP_PLACEHOLDER.OBJECT and hasattr(shape, "image") and getattr(shape, "image"):
                    pic = process_picture(config, shape, slide_id)
                    if pic:
                        initial_elements.append(pic)
            except:
                pass # Ignore shapes that are not text, pic, table, or recognized object

    # Refine initial elements: convert math paragraphs, merge code blocks
    processed_elements = _refine_elements(initial_elements, slide_id)
            
    return processed_elements


def parse(config: ConversionConfig, prs: Presentation) -> ParsedPresentation:
    result = ParsedPresentation(slides=[])

    for idx, slide in enumerate(tqdm(prs.slides, desc='Converting slides')):
        if config.page is not None and idx + 1 != config.page:
            continue
        shapes = []
        try:
            shapes = sorted(ungroup_shapes(slide.shapes), key=attrgetter('top', 'left'))
        except:
            logger.warning('Bad shapes encountered in this slide. Please check or remove them and try again.')
            logger.warning('shapes:')
            try:
                for sp in slide.shapes:
                    logger.warning(sp.shape_type)
                    logger.warning(sp.top, sp.left, sp.width, sp.height)
            except:
                logger.warning('failed to print all bad shapes.')

        if not config.try_multi_column:
            result_slide = GeneralSlide(elements=process_shapes(config, shapes, idx + 1))
        else:
            multi_column_slide = get_multi_column_slide_if_present(
                prs, slide, partial(process_shapes, config=config, slide_id=idx + 1))
            if multi_column_slide:
                result_slide = multi_column_slide
            else:
                result_slide = GeneralSlide(elements=process_shapes(config, shapes, idx + 1))

        if not config.disable_notes and slide.has_notes_slide:
            text = slide.notes_slide.notes_text_frame.text
            if text:
                result_slide.notes.append(text)

        result.slides.append(result_slide)

    return result
